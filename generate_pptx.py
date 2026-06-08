"""
generate_pptx.py  —  Generate QR Code Generator Bot Presentation
=================================================================
Generates QR_code_generator/project_demo.pptx with 13 slides,
dark corporate theme, architecture diagrams, and speaker notes.

Usage:  python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ACCENT_ORANGE = RGBColor(0xFF, 0xA5, 0x00)
TEXT_WHITE    = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_GRAY     = RGBColor(0xA0, 0xA0, 0xB0)
TEXT_DIM      = RGBColor(0x70, 0x70, 0x80)

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "QR_code_generator", "project_demo.pptx"
)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            a.set('val', str(int(alpha * 1000)))
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=TEXT_WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    add_rect(slide, left, top, width, Pt(3), color)

def add_section_header(slide, number, title, subtitle=None):
    add_bg(slide)
    # Accent bar at top
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT_BLUE)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.5), Inches(1.2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number).zfill(2)
    p.font.size = Pt(36)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    # Title
    add_textbox(slide, Inches(2.4), Inches(1.5), Inches(10), Inches(1),
                title, font_size=36, bold=True, color=TEXT_WHITE)
    add_accent_line(slide, Inches(2.4), Inches(2.6), Inches(3), ACCENT_GREEN)
    if subtitle:
        add_textbox(slide, Inches(2.4), Inches(2.9), Inches(10), Inches(0.8),
                    subtitle, font_size=18, color=TEXT_GRAY)

def add_card(slide, left, top, width, height, title, items, title_color=ACCENT_BLUE):
    add_rect(slide, left, top, width, height, BG_CARD)
    add_accent_line(slide, left, top, width, title_color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                title, font_size=16, bold=True, color=title_color)
    add_bullet_slide(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                     height - Inches(0.9), items, font_size=13, color=TEXT_GRAY, spacing=Pt(4))


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Pt(5), ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
            "QR Code Generator Bot", font_size=48, bold=True, color=TEXT_WHITE,
            alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(3.3), Inches(3.3), ACCENT_GREEN)
add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
            "n8n + Telegram + Cloudflare Tunnel + QRServer API",
            font_size=22, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
            "Nitin Kumar  •  2026", font_size=16, color=TEXT_DIM,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
            "https://github.com/nitinkumar30/n8n-projects", font_size=12,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Welcome to the QR Code Generator Bot presentation. This project demonstrates a production-ready Telegram bot built on n8n with free HTTPS via Cloudflare Tunnel."

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "The Problem", "Why existing solutions fall short")
# Pain points
pains = [
    "Telegram bots require HTTPS — localhost and HTTP are rejected",
    "ngrok provides HTTPS but URLs change every session",
    "Manual reconfiguration needed every time the tunnel restarts",
    "Test workflow mode stops listening after one execution",
    "No simple way to keep a bot running continuously on a local machine"
]
add_bullet_slide(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
                 pains, font_size=18, color=TEXT_WHITE, bullet_color=ACCENT_ORANGE)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "The core challenge: Telegram's API mandates HTTPS webhooks. Localhost is blocked. Free ngrok tunnels change URLs constantly, requiring manual environment variable updates. Test mode in n8n only runs once. This project solves all of this with a zero-touch setup."

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "The Solution", "Cloudflare Tunnel + Python automation")
items = [
    "Replace ngrok with Cloudflare Tunnel for stable, free HTTPS",
    "Python auto-start script manages tunnel, env vars, and n8n launch",
    "Workflow runs in ACTIVE mode for continuous listening",
    "URL validation prevents invalid requests from reaching QR API",
    "Clean, styled QR codes with configurable size and colors"
]
add_bullet_slide(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
                 items, font_size=18, color=TEXT_WHITE, bullet_color=ACCENT_GREEN)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Cloudflare Tunnel gives us a stable HTTPS URL that doesn't change every session. The Python script extracts the tunnel URL automatically and sets the required n8n environment variables. The workflow stays in ACTIVE mode, listening indefinitely."

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Architecture Overview", "End-to-end data flow")
# Flow boxes
flow = [
    ("Telegram\nUser", ACCENT_BLUE, Inches(0.8)),
    ("Telegram\nServers", ACCENT_PURPLE, Inches(3.0)),
    ("Cloudflare\nTunnel", ACCENT_ORANGE, Inches(5.2)),
    ("n8n\nWebhook", ACCENT_GREEN, Inches(7.4)),
    ("QRServer\nAPI", ACCENT_BLUE, Inches(9.6)),
]
for label, color, left in flow:
    add_rect(slide, left, Inches(3.2), Inches(1.8), Inches(1.4), color)
    add_textbox(slide, left, Inches(3.5), Inches(1.8), Inches(1.2),
                label, font_size=14, bold=True, color=TEXT_WHITE,
                alignment=PP_ALIGN.CENTER)
    # Arrow
    if left < 9.6:
        add_textbox(slide, left + Inches(1.8), Inches(3.6), Inches(1.2), Inches(0.5),
                    "→", font_size=28, bold=True, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1),
            "Telegram User → Telegram Servers → Cloudflare Tunnel → n8n Webhook → QRServer API → Response",
            font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "The data flow: User sends a message in Telegram → Telegram servers receive it → Cloudflare Tunnel forwards the HTTPS webhook to local n8n → n8n validates and calls QRServer API → QR image is sent back through the same chain. All HTTPS is handled by Cloudflare, so Telegram is happy."

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Tech Stack
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "Technology Stack", "Tools and services powering the bot")
stack = [
    ("n8n", "Workflow automation engine with visual builder. Handles webhooks, logic, and API calls."),
    ("Telegram Bot API", "Messaging platform. Receive messages and send QR images via bot token."),
    ("Cloudflare Tunnel", "Free HTTPS tunnel. Exposes localhost:5678 as a public HTTPS URL."),
    ("QRServer API", "QR code generation. Styled outputs with configurable size, color, and ECC."),
    ("Python", "Auto-start script. Manages tunnel lifecycle, env vars, and n8n process."),
]
y = Inches(3.4)
for title, desc in stack:
    add_textbox(slide, Inches(1.2), y, Inches(2.5), Inches(0.5),
                title, font_size=18, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, Inches(3.8), y, Inches(8.5), Inches(0.5),
                desc, font_size=15, color=TEXT_GRAY)
    y += Inches(0.7)

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Five technologies work together: n8n orchestrates the workflow, Telegram provides the chat interface, Cloudflare Tunnel handles HTTPS, QRServer generates the codes, and Python ties it all together with automation."

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Workflow: Node Map
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Workflow Node Map", "4 nodes, one clean flow")
# Node cards
nodes = [
    ("1. Telegram Trigger", "Listens for incoming messages via webhook. Configured for 'message' events with bot token credential."),
    ("2. IF Node (Validation)", "Validates URL using regex ^(https?:\\/\\/)[^\\s]+$. Blocks invalid input before it reaches the QR API."),
    ("3. HTTP Request", "Calls api.qrserver.com with styled parameters: size=350x350, ECC=M, custom color, white background, 4px quiet zone."),
    ("4. Telegram Send Photo", "Sends the QR image back to the user with caption. Uses binary property 'data' from HTTP response."),
]
y = Inches(3.3)
for title, desc in nodes:
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), BG_CARD)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(3.5), Inches(0.4),
                title, font_size=15, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, Inches(4.6), y + Inches(0.05), Inches(7.5), Inches(0.7),
                desc, font_size=13, color=TEXT_GRAY)
    y += Inches(0.95)

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Four nodes: Telegram Trigger listens, IF node validates the URL is well-formed, HTTP Request generates the QR with styled parameters, and Telegram Send Photo delivers the image. Simple but effective."

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Inputs
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "Inputs & Validation", "What users send, what the bot accepts")
add_textbox(slide, Inches(1), Inches(3.5), Inches(5), Inches(0.5),
            "✅ Accepted", font_size=20, bold=True, color=ACCENT_GREEN)
valid = [
    "https://github.com",
    "http://example.com/page?q=1",
    "https://sub.domain.org/path/file.html",
    "ftp://files.server.com (also accepted)"
]
add_bullet_slide(slide, Inches(1), Inches(4.1), Inches(5), Inches(2.5),
                 valid, font_size=16, color=TEXT_WHITE, bullet_color=ACCENT_GREEN)

add_textbox(slide, Inches(6.5), Inches(3.5), Inches(5), Inches(0.5),
            "❌ Rejected", font_size=20, bold=True, color=RGBColor(0xFF, 0x55, 0x55))
invalid = [
    "hello (plain text)",
    "192.168.1.1 (not a URL)",
    "  https://example.com (leading space)",
    "just random words"
]
add_bullet_slide(slide, Inches(6.5), Inches(4.1), Inches(5), Inches(2.5),
                 invalid, font_size=16, color=TEXT_WHITE,
                 bullet_color=RGBColor(0xFF, 0x55, 0x55))

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "The bot accepts any valid URL starting with http://, https://, or ftp://. Everything else — plain text, IP addresses, malformed input — gets rejected with a clear error message. This prevents wasteful API calls."

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Outputs
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "Outputs & User Experience", "What the user receives")
outputs = [
    "Styled QR code image (350×350, custom blue color, white background)",
    "Error message for invalid input: 'Please send valid URL'",
    "Consistent caption: 'Here is your QR code'",
    "Response in under 2 seconds (Cloudflare + QRServer latency)",
    "No rate limiting on the bot side (subject to Telegram's limits)"
]
add_bullet_slide(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
                 outputs, font_size=18, color=TEXT_WHITE, bullet_color=ACCENT_BLUE)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Users get a clean 350x350 QR code in custom blue with proper quiet zone margins. Invalid input gets a friendly but clear error. Response time is typically under 2 seconds."

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Benefits
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 8, "Key Benefits", "Why this approach wins")
# Benefit cards
benefits = [
    ("Zero Cost", "Cloudflare Tunnel is free. QRServer is free. n8n is free. No paid subscriptions needed."),
    ("Always-On HTTPS", "Cloudflare provides stable HTTPS that doesn't change URLs every session like ngrok."),
    ("One-Click Start", "Python script handles everything: tunnel, env vars, n8n launch. Run one command, done."),
    ("Continuous Operation", "ACTIVE mode means the bot listens indefinitely. No manual restart needed."),
    ("Clean QR Codes", "Styled with custom colors, ECC correction, and proper quiet zone for better scanning."),
]
y = Inches(3.3)
for title, desc in benefits:
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), BG_CARD)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.4),
                title, font_size=15, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, Inches(3.6), y + Inches(0.05), Inches(8.5), Inches(0.6),
                desc, font_size=13, color=TEXT_GRAY)
    y += Inches(0.82)

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Five key benefits: completely free, always-on HTTPS, one-command startup, continuous listening without babysitting, and professional-looking QR codes with custom styling."

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Security
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 9, "Security & Best Practices", "Running safely")
sec_items = [
    "HTTPS end-to-end: Telegram → Cloudflare → n8n (no plaintext webhooks)",
    "URL validation blocks malicious or malformed input before processing",
    "No sensitive data stored — QR codes are generated on-the-fly and discarded",
    "Bot token stored in n8n credential system (encrypted at rest)",
    "Local machine only — PC must be ON and connected to internet",
    "No authentication on the QRServer API (public service, use responsibly)"
]
add_bullet_slide(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
                 sec_items, font_size=17, color=TEXT_WHITE, bullet_color=ACCENT_BLUE)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Security considerations: HTTPS is end-to-end via Cloudflare. URL validation prevents abuse. No data persistence. Bot token is encrypted in n8n's credential vault. The only risk is the PC being off or losing internet, which stops the bot."

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Demo
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 10, "Live Demo", "See it in action")
demo_items = [
    "Send: https://github.com",
    "Bot replies with QR code for github.com",
    "",
    "Send: hello",
    "Bot replies: 'Please send valid URL'",
    "",
    "Send: https://www.google.com/search?q=n8n",
    "Bot replies with QR for the search URL",
    "",
    "All within 2 seconds per request",
]
add_bullet_slide(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
                 demo_items, font_size=18, color=TEXT_WHITE, bullet_color=ACCENT_GREEN)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Live demo: Three example interactions showing valid URL, invalid input, and complex URL. Each responds within 2 seconds. The bot handles any properly formatted URL."

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Roadmap
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 11, "Future Roadmap", "What's coming next")
roadmap = [
    ("Custom QR Themes", "Allow users to select colors, add logos, or choose patterns via Telegram commands"),
    ("QR History", "Save generated QR codes to a database and let users browse past codes"),
    ("VPS Deployment", "Deploy to a cloud VPS for 24/7 uptime instead of depending on local PC"),
    ("UPI QR Generator", "Generate UPI payment QR codes for Indian users with amount support"),
    ("Multi-Command Bot", "Add more commands: /help, /history, /theme, /stats"),
    ("Analytics Dashboard", "Track usage metrics: number of QR codes generated, popular URLs, errors"),
]
y = Inches(3.3)
for title, desc in roadmap:
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), BG_CARD)
    add_textbox(slide, Inches(1.1), y + Inches(0.02), Inches(3), Inches(0.35),
                title, font_size=14, bold=True, color=ACCENT_ORANGE)
    add_textbox(slide, Inches(4.2), y + Inches(0.02), Inches(8), Inches(0.6),
                desc, font_size=12, color=TEXT_GRAY)
    y += Inches(0.78)

add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "The roadmap covers six areas: custom themes, QR history with database, 24/7 VPS deployment, UPI payment QR support, multi-command expansion, and an analytics dashboard. Priority is VPS deployment for always-on availability."

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Pt(5), ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
            "Thank You", font_size=48, bold=True, color=TEXT_WHITE,
            alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(3.3), Inches(3.3), ACCENT_GREEN)
add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
            "Questions? Feedback? Let's connect.",
            font_size=22, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
            "Nitin Kumar", font_size=20, bold=True, color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
            "linkedin.com/in/nitin30kumar  •  github.com/nitinkumar30",
            font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.5),
            "github.com/nitinkumar30/n8n-projects",
            font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_notes_slide = slide.notes_slide
notes_tf = add_notes_slide.notes_text_frame
notes_tf.text = "Thank you for your time. The full source code and documentation is available on GitHub. Connect with me on LinkedIn or explore the repository for more n8n projects."

# ── Save ────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"Presentation saved: {OUTPUT_PATH}")
